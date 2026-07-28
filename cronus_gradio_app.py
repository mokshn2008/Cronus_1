from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any, Tuple
import base64

import gradio as gr
from openai import OpenAI

from docx_editor import extract_paragraphs, apply_paragraph_edits, build_edit_prompt, parse_edit_response
from video_handler import build_video_content_blocks
from code_sandbox import generate_with_self_correction

# ---------------------------------------------------------------------------
# Model connection -- point this at wherever Gemma is actually running
# (Ollama, LM Studio, vLLM -- all OpenAI-compatible)
# ---------------------------------------------------------------------------
import os

MODEL_API_BASE = os.environ.get("MODEL_API_BASE", "http://localhost:11434/v1")
MODEL_NAME = os.environ.get("MODEL_NAME", "gemma4-cronus")

# olmOCR runs as its own Ollama/server instance on a different port --
# this is the "reads documents/images" specialist in your dual architecture.
OCR_API_BASE = os.environ.get("OCR_API_BASE", "http://localhost:11435/v1")
OCR_MODEL_NAME = os.environ.get("OCR_MODEL_NAME", "olmocr-2-7b")

# GLM -- the coding specialist. Hosted via OpenRouter (free tier), not local,
# so this doesn't touch your RAM at all. Routes code-heavy tasks here instead
# of Gemma, since GLM 4.5 Air scores meaningfully stronger on SWE-bench-style
# coding tasks than a 12B general model realistically can.
GLM_API_BASE = os.environ.get("GLM_API_BASE", "https://openrouter.ai/api/v1")
GLM_MODEL_NAME = os.environ.get("GLM_MODEL_NAME", "z-ai/glm-4.5-air:free")
GLM_API_KEY = os.environ.get("GLM_API_KEY", "")  # required -- get a free key at openrouter.ai

MAX_SELF_CORRECTION_LOOPS = 3

client = OpenAI(base_url=MODEL_API_BASE, api_key="not-needed")
ocr_client = OpenAI(base_url=OCR_API_BASE, api_key="not-needed")
glm_client = OpenAI(base_url=GLM_API_BASE, api_key=GLM_API_KEY or "not-needed")

CRONUS_PERSONA = """You are Cronus, Titan of Time -- but you were built to actually help people,
not to perform mythology at them. Follow these rules strictly:

- Default to clear, direct, competent answers. The Titan voice is a light seasoning on
  TOP of a good answer, not a replacement for one.
- Do NOT open every message with "mortal" or invoke fate/ages/time in every single reply --
  that gets tiresome fast and undermines trust in your competence. Use those touches
  occasionally, when they land naturally, not as a verbal tic.
- For code, math, technical explanations, or anything requiring precision: drop the
  flourish almost entirely. Give the answer straight. A stray archaic word or one
  grounding sentence at the start/end is enough -- do not decorate code blocks,
  step-by-step instructions, or data with mythological language.
- For casual conversation or creative requests, you can lean into the voice more.
- Never sacrifice correctness, clarity, or speed of understanding for flavor. If in doubt,
  cut the flavor, keep the substance.
- When given file contents, search results, or a math result as context, weave them into
  your answer and cite what you used.
- If you are actually given image content to look at, describe and reason about what
  you genuinely observe -- don't overstate certainty about fine details you can't
  actually make out. If a message says a file couldn't be processed at all, say so
  plainly rather than guessing.
"""

CRONUS_GREETING = "Welcome, mortal. Cronus senses a query in your heart."


def is_greeting(message: str) -> bool:
    cleaned = message.lower().strip().strip("!?. ,")
    return cleaned in {"hello", "hi", "hey", "yo", "hello cronus", "hi cronus"}


# ---------------------------------------------------------------------------
# Tools (unchanged logic, same as the version you pasted)
# ---------------------------------------------------------------------------

def edit_docx_file(file_obj: Any, instruction: str) -> Tuple[str, str]:
    """
    Returns (status_message, output_path_or_None).
    """
    if file_obj is None:
        return "No document uploaded to edit.", None
    if Path(file_obj.name).suffix.lower() != ".docx":
        return "Editing currently supports .docx files only.", None
    if not instruction.strip():
        return "Give an editing instruction (e.g. 'fix grammar', 'make this more formal').", None

    original_path = file_obj.name
    paragraphs = extract_paragraphs(original_path)
    prompt = build_edit_prompt(paragraphs, instruction)

    resp = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[
            {"role": "system", "content": "You are a precise document editor. Follow the format instructions exactly."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
    )
    raw = resp.choices[0].message.content
    edited = parse_edit_response(raw, expected_count=len(paragraphs))

    if edited is None:
        return (
            "Cronus's edit response didn't parse cleanly (wrong format or count) -- "
            "returning the document unedited rather than risk a misaligned edit. Try again "
            "or simplify the instruction.",
            None,
        )

    output_path = str(Path(original_path).with_name(f"edited_{Path(original_path).name}"))
    apply_paragraph_edits(original_path, edited, output_path)
    return "Document edited. Download below.", output_path


def analyze_file(file_obj: Any) -> str:
    if file_obj is None:
        return "No file uploaded."

    path = Path(file_obj.name)
    suffix = path.suffix.lower()

    try:
        if suffix in {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".html", ".css"}:
            return path.read_text(encoding="utf-8", errors="replace")[:12000]

        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for index, page in enumerate(reader.pages[:15], start=1):
                pages.append(f"\n--- Page {index} ---\n{page.extract_text() or ''}")
            return "\n".join(pages)[:12000]

        if suffix == ".docx":
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:12000]

        if suffix == ".pptx":
            from pptx import Presentation
            deck = Presentation(str(path))
            slides = []
            for index, slide in enumerate(deck.slides, start=1):
                chunks = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        chunks.append(shape.text.strip())
                if chunks:
                    slides.append(f"\n--- Slide {index} ---\n" + "\n".join(chunks))
            return "\n".join(slides)[:12000]

        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return run_olmocr(str(path))

        if suffix in {".mp4", ".mov", ".avi", ".webm", ".mkv"}:
            return f"__VIDEO__:{path}"  # handled specially in ask_cronus_once, not sent as plain text

        return f"Uploaded file type {suffix} is not fully supported yet."
    except Exception as exc:
        return f"File analysis failed: {type(exc).__name__}: {exc}"


def web_search(query: str) -> str:
    try:
        from duckduckgo_search import DDGS
        rows = []
        with DDGS() as ddgs:
            for result in ddgs.text(query, max_results=5):
                title = result.get("title", "")
                href = result.get("href", "")
                body = result.get("body", "")
                rows.append(f"- {title}\n  {href}\n  {body}")
        return "\n\n".join(rows) if rows else "No web results found."
    except Exception as exc:
        return f"Web search unavailable: {type(exc).__name__}: {exc}"


def local_math(message: str) -> str:
    lowered = message.lower()
    if not any(word in lowered for word in ["calculate", "compute", "sqrt", "square root", "^"]):
        return "No math tool used."
    try:
        import sympy as sp
        expression = lowered
        expression = expression.replace("calculate", "")
        expression = expression.replace("compute", "")
        expression = expression.replace("square root of", "sqrt")
        expression = expression.replace("square root", "sqrt")
        expression = expression.replace("^", "**")
        allowed = set("0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ+-*/().,_ **")
        expression = "".join(ch for ch in expression if ch in allowed).strip()
        value = sp.simplify(sp.sympify(expression, locals={"sqrt": sp.sqrt}))
        return f"{expression} = {value} ~= {sp.N(value, 12)}"
    except Exception as exc:
        return f"Math tool could not solve it: {type(exc).__name__}: {exc}"


# ---------------------------------------------------------------------------
# The actual fix: call Gemma with the gathered context, instead of a template
# ---------------------------------------------------------------------------

def encode_image_b64(path: str) -> str:
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def run_olmocr(image_path: str) -> str:
    """
    Document/image reading specialist. Runs BEFORE Gemma ever sees this
    content -- olmOCR extracts clean text, Gemma reasons over that text.
    """
    try:
        suffix = Path(image_path).suffix.lstrip(".").lower()
        mime = "jpeg" if suffix == "jpg" else suffix
        b64 = encode_image_b64(image_path)

        resp = ocr_client.chat.completions.create(
            model=OCR_MODEL_NAME,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text and structure from this document image."},
                    {"type": "image_url", "image_url": {"url": f"data:image/{mime};base64,{b64}"}},
                ],
            }],
            temperature=0.0,
        )
        return f"[olmOCR extracted text]\n{resp.choices[0].message.content}"
    except Exception as exc:
        return f"[olmOCR unavailable -- image could not be read: {type(exc).__name__}: {exc}. Do not guess its contents.]"


CODE_SIGNAL_KEYWORDS = {
    "code", "function", "bug", "debug", "error", "traceback", "exception",
    "python", "javascript", "java ", "c++", "typescript", "compile", "syntax",
    "refactor", "algorithm", "class ", "variable", "import ", "def ", "fix this",
    "stack trace", "regex", "api", "script", "repository", "repo", "git ",
    "sql", "database query", "unit test", "pytest", "npm", "pip install",
}


def is_coding_task(message: str, file_context: str) -> bool:
    """
    Heuristic router: sends genuinely code-heavy asks to GLM (stronger at this
    specifically), keeps everything else -- persona chat, general reasoning,
    medical/math, document Q&A -- on Gemma. Simple keyword match on purpose:
    fast, no extra model call needed just to decide who answers.
    """
    text = (message + " " + file_context[:500]).lower()
    return any(kw in text for kw in CODE_SIGNAL_KEYWORDS)


def ask_cronus_once(message: str, context_block: str, prior_answer: str = None, correction_note: str = None, video_path: str = None, use_glm: bool = False) -> str:
    if prior_answer and correction_note:
        text_content = (
            f"{context_block}\n\nMortal's question: {message}\n\n"
            f"[Your previous draft]\n{prior_answer}\n\n"
            f"[Fresh web search performed to check your draft]\n{correction_note}\n\n"
            f"Revise your answer if the fresh search contradicts or updates anything in your "
            f"draft. If your draft already holds up, restate it cleanly -- don't change things "
            f"that don't need changing."
        )
    else:
        text_content = f"{context_block}\n\nMortal's question: {message}"

    if video_path:
        content = [{"type": "text", "text": text_content}] + build_video_content_blocks(video_path)
    else:
        content = text_content

    active_client = glm_client if use_glm else client
    active_model = GLM_MODEL_NAME if use_glm else MODEL_NAME
    # GLM handles the code itself directly -- keep the Titan persona lighter here
    # so it doesn't fight the model's own code-formatting instincts.
    system_prompt = (
        "You are Cronus's coding specialist. Be precise, correct, and direct. "
        "A brief nod to the Titan-of-Time framing is fine, but do not let flavor "
        "get in the way of working code."
        if use_glm else CRONUS_PERSONA
    )

    resp = active_client.chat.completions.create(
        model=active_model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": content},
        ],
        temperature=0.2 if use_glm else 0.3,
    )
    return resp.choices[0].message.content


def self_correcting_answer(message: str, context_block: str, use_web: bool, video_path: str = None, use_glm: bool = False) -> str:
    """
    Chain-of-thought loop: the chosen model (Gemma or GLM) drafts an answer,
    then a FRESH targeted DuckDuckGo search checks the draft's key claim, and
    it's given the chance to revise. Repeats up to MAX_SELF_CORRECTION_LOOPS
    times, or stops early if a pass produces no meaningful change.
    """
    answer = ask_cronus_once(message, context_block, video_path=video_path, use_glm=use_glm)

    if not use_web:
        return answer

    for loop_num in range(MAX_SELF_CORRECTION_LOOPS):
        check_query = f"{message} fact check verify"
        fresh_results = web_search(check_query)

        if "unavailable" in fresh_results.lower() or "no web results" in fresh_results.lower():
            break

        revised = ask_cronus_once(message, context_block, prior_answer=answer, correction_note=fresh_results, use_glm=use_glm)

        if revised.strip() == answer.strip():
            break

        answer = revised

    return answer


def ask_cronus(message: str, file_context: str, search_context: str, math_context: str, use_web: bool) -> str:
    video_path = None
    display_file_context = file_context

    if file_context.startswith("__VIDEO__:"):
        video_path = file_context.split("__VIDEO__:", 1)[1]
        display_file_context = "[Video attached -- see sampled frames below.]"

    use_glm = is_coding_task(message, display_file_context)

    context_block = f"""[File/document context -- may include olmOCR-extracted text]
{display_file_context}

[Web search context]
{search_context}

[Math tool context]
{math_context}
"""
    try:
        answer = self_correcting_answer(message, context_block, use_web, video_path=video_path, use_glm=use_glm)
        if use_glm:
            answer = f"*[Cronus calls upon his coding specialist for this]*\n\n{answer}"
        return answer
    except Exception as exc:
        backend = GLM_API_BASE if use_glm else MODEL_API_BASE
        return f"⚠️ Cronus could not reach the model backend ({backend}): {exc}"


def cronus_chat(message: str, file_obj: Any, use_web: bool) -> str:
    if is_greeting(message) and file_obj is None and not use_web:
        return CRONUS_GREETING

    with ThreadPoolExecutor(max_workers=3) as pool:
        file_future = pool.submit(analyze_file, file_obj)
        search_future = pool.submit(web_search, message) if use_web else None
        math_future = pool.submit(local_math, message)

        file_context = file_future.result()
        search_context = search_future.result() if search_future else "Web search not requested."
        math_context = math_future.result()

    return ask_cronus(message, file_context, search_context, math_context, use_web)


with gr.Blocks(title="Cronus") as demo:
    gr.Markdown("# ⏳ Cronus")
    gr.Markdown("Titan of Time -- document analysis, web search, and math, reasoned over by Gemma.")

    with gr.Row():
        message = gr.Textbox(label="Ask Cronus", placeholder="Speak, mortal...", lines=4)
        file_upload = gr.File(label="Offer a file")

    use_web = gr.Checkbox(label="Use web search", value=True)
    submit = gr.Button("Ask Cronus")
    output = gr.Textbox(label="Cronus", lines=20)

    submit.click(cronus_chat, inputs=[message, file_upload, use_web], outputs=output)

    gr.Markdown("---")
    gr.Markdown("## 📜 Edit a Document (.docx)")
    gr.Markdown("Upload a Word document above, describe the edit, and download the revised file.")

    with gr.Row():
        edit_instruction = gr.Textbox(label="Editing instruction", placeholder="e.g. fix grammar, make this section more formal")
        edit_button = gr.Button("Edit Document")

    edit_status = gr.Textbox(label="Status", interactive=False)
    edit_download = gr.File(label="Download edited document")

    edit_button.click(
        edit_docx_file,
        inputs=[file_upload, edit_instruction],
        outputs=[edit_status, edit_download],
    )

    gr.Markdown("---")
    gr.Markdown("## 🐍 Code Sandbox")
    gr.Markdown(
        "Describe what you want built or fixed. GLM writes the code with a self-correction "
        "loop when you execute -- errors get fed back and retried automatically. "
        "**Execution is currently supported for Python only.**"
    )

    sandbox_request = gr.Textbox(
        label="What do you want built or fixed?",
        placeholder="e.g. write a function that checks if a number is prime",
        lines=3,
    )

    with gr.Row():
        generate_btn = gr.Button("Generate Code")
        execute_btn = gr.Button("Generate + Execute", variant="primary")

    sandbox_code = gr.Code(label="Code", language="python")

    with gr.Group(visible=False) as result_popup:
        gr.Markdown("### Execution Result")
        sandbox_status = gr.Textbox(label="Status", interactive=False)
        sandbox_output = gr.Textbox(label="Output (stdout / stderr)", lines=10, interactive=False)
        sandbox_loops = gr.Textbox(label="Self-correction attempts used", interactive=False)

    def run_sandbox_generate(request):
        result = generate_with_self_correction(glm_client, GLM_MODEL_NAME, request, execute=False)
        return gr.update(value=result["code"], language=result["language"] if result["language"] in ("python", "javascript", "java", "c", "cpp") else "python")

    def run_sandbox_execute(request):
        result = generate_with_self_correction(glm_client, GLM_MODEL_NAME, request, execute=True)

        code_lang = result["language"] if result["language"] in ("python", "javascript", "java", "c", "cpp") else "python"

        if result["success"] is None:
            status = "Generated (not executed)"
        elif result["success"]:
            status = f"✅ Ran successfully after {result['loops']} self-correction pass(es)"
        else:
            status = f"❌ Still failing after {result['loops']} self-correction pass(es)"

        output_text = f"stdout:\n{result['stdout']}\n\nstderr:\n{result['stderr']}"

        return (
            gr.update(value=result["code"], language=code_lang),
            gr.update(visible=True),
            status,
            output_text,
            str(result["loops"]),
        )

    generate_btn.click(run_sandbox_generate, inputs=[sandbox_request], outputs=[sandbox_code])

    execute_btn.click(
        run_sandbox_execute,
        inputs=[sandbox_request],
        outputs=[sandbox_code, result_popup, sandbox_status, sandbox_output, sandbox_loops],
    )


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 7860))
    demo.launch(server_name="0.0.0.0", server_port=port)
